import io

import requests
import streamlit as st
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.util import Inches

from services.api_client import generate_executive_report_narrative, list_workspaces

st.title("PowerPoint Generator")

if "workspace_id" not in st.session_state:
    st.warning("Select or create a workspace on the Home page first.")
    st.stop()

workspace_id = st.session_state["workspace_id"]
workspace = next((w for w in list_workspaces() if w["id"] == workspace_id), None)

swot = st.session_state.get("swot_result")
roadmap = st.session_state.get("roadmap_result")
maturity = st.session_state.get("maturity_result")
technology = st.session_state.get("technology_result")
kpi_df = st.session_state.get("kpi_results")
kpi = kpi_df.to_dict(orient="records") if kpi_df is not None else None
pain_points = st.session_state.get("interview_findings")

st.subheader("Available findings")
status = [
    ("Pain Points (Interview Analyzer)", pain_points),
    ("SWOT", swot),
    ("Digital Maturity", maturity),
    ("Roadmap", roadmap),
    ("Technology Recommendations", technology),
    ("KPI Data", kpi),
]
for label, value in status:
    st.write(f"{'✅' if value else '⬜'} {label}")

# Reuse Module 11's narrative endpoint and its result if the Executive
# Report page already generated one in this session — no reason to pay
# for a second Claude call producing the same content.
if "executive_narrative" not in st.session_state:
    if not any(value for _, value in status):
        st.info(
            "Nothing generated yet — visit the Interview Analyzer, SWOT, Roadmap, "
            "Digital Maturity, Technology, or KPI pages first, then come back here."
        )
        st.stop()

    if st.button("Generate Slide Content"):
        try:
            with st.spinner("Synthesizing findings into narrative content..."):
                narrative = generate_executive_report_narrative(
                    workspace_id,
                    {
                        "swot": swot,
                        "roadmap": roadmap,
                        "maturity": maturity,
                        "technology": technology,
                        "kpi": kpi,
                    },
                )
        except requests.exceptions.RequestException as exc:
            st.error(f"Could not generate slide content. ({exc})")
            st.stop()
        st.session_state["executive_narrative"] = narrative
        st.rerun()
    st.stop()

narrative = st.session_state["executive_narrative"]
st.success("Narrative content ready (reused from Executive Report if already generated).")

MATURITY_CATEGORIES = [
    "leadership",
    "operations",
    "technology",
    "data",
    "supply_chain",
    "automation",
    "sustainability",
    "cybersecurity",
    "workforce",
]


def build_pptx() -> bytes:
    prs = Presentation()
    title_layout = prs.slide_layouts[0]
    content_layout = prs.slide_layouts[1]
    chart_layout = prs.slide_layouts[5]  # Title Only — leaves room for a chart

    def add_bullet_slide(title: str, lines: list[str] | None):
        slide = prs.slides.add_slide(content_layout)
        slide.shapes.title.text = title
        lines = lines or ["Not available for this engagement."]
        body = slide.placeholders[1].text_frame
        # Without these, a slide with many/long bullets just overflows the
        # placeholder silently instead of shrinking to fit — confirmed
        # this combination (word_wrap + TEXT_TO_FIT_SHAPE) works together.
        body.word_wrap = True
        body.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        body.text = lines[0]
        for line in lines[1:]:
            paragraph = body.add_paragraph()
            paragraph.text = line

    def add_text_slide(title: str, text: str | None):
        add_bullet_slide(title, [text] if text else None)

    def add_chart_slide(title: str, categories: list[str], values: list[float]):
        # A real, native PowerPoint chart object (editable in PowerPoint
        # itself) rather than a static image — reads far better than a
        # slide of "Category: score" bullet lines for numeric data.
        slide = prs.slides.add_slide(chart_layout)
        slide.shapes.title.text = title
        chart_data = CategoryChartData()
        chart_data.categories = categories
        chart_data.add_series("Score", values)
        graphic_frame = slide.shapes.add_chart(
            XL_CHART_TYPE.COLUMN_CLUSTERED,
            Inches(0.5), Inches(1.5), Inches(9), Inches(5),
            chart_data,
        )
        graphic_frame.chart.has_legend = False
        return slide

    # 1. Title
    title_slide = prs.slides.add_slide(title_layout)
    client_name = workspace["client_name"] if workspace else "Client"
    industry = workspace["industry"] if workspace else ""
    title_slide.shapes.title.text = client_name
    title_slide.placeholders[1].text = f"Digital Transformation Assessment — {industry}"

    # 2. Executive Summary
    add_text_slide("Executive Summary", narrative.get("executive_summary"))

    # 3. Current State
    add_text_slide("Current State", narrative.get("current_situation"))

    # 4. Pain Points
    pain_lines = [
        f"{f['pain_point']} ({f['severity']}): {f['recommendation']}" for f in (pain_points or [])
    ]
    add_bullet_slide("Pain Points", pain_lines)

    # 5. Operational Analysis
    if kpi:
        categories = [row.get("machine") or "?" for row in kpi]
        oee_values = [round((row.get("oee") or 0) * 100, 1) for row in kpi]
        add_chart_slide("Operational Analysis — OEE by Machine (%)", categories, oee_values)
    else:
        add_bullet_slide("Operational Analysis", None)

    # 6. Digital Maturity
    if maturity:
        categories = [c.replace("_", " ").title() for c in MATURITY_CATEGORIES]
        scores = [(maturity.get(c) or {}).get("score", 0) for c in MATURITY_CATEGORIES]
        add_chart_slide(f"Digital Maturity — Overall {maturity.get('overall')}/5", categories, scores)
    else:
        add_bullet_slide("Digital Maturity", None)

    # 7. Roadmap
    roadmap_lines = []
    if roadmap:
        for horizon, label in [
            ("quick_wins", "Quick Wins"),
            ("medium_term", "Medium-Term"),
            ("long_term", "Long-Term"),
        ]:
            for item in roadmap.get(horizon, []):
                roadmap_lines.append(f"[{label}] {item.get('initiative')}")
    add_bullet_slide("Transformation Roadmap", roadmap_lines)

    # 8. Technology Recommendations
    tech_lines = []
    if technology:
        for rec in technology.get("recommendations", []):
            tech_lines.append(
                f"{rec.get('recommendation')} via {rec.get('platform')} "
                f"(ROI: {rec.get('expected_return')})"
            )
    add_bullet_slide("Technology Recommendations", tech_lines)

    # 9. Financial Benefits
    add_text_slide("Financial Benefits", narrative.get("financial_impact"))

    # 10. Conclusion
    add_bullet_slide("Conclusion", narrative.get("key_findings"))

    # 11. Future Steps
    add_bullet_slide("Future Steps", narrative.get("next_steps"))

    buffer = io.BytesIO()
    prs.save(buffer)
    return buffer.getvalue()


st.divider()
st.subheader("Preview")
for heading in [
    "Title",
    "Executive Summary",
    "Current State",
    "Pain Points",
    "Operational Analysis",
    "Digital Maturity",
    "Transformation Roadmap",
    "Technology Recommendations",
    "Financial Benefits",
    "Conclusion",
    "Future Steps",
]:
    st.write(f"- {heading}")

st.download_button(
    "Download PowerPoint",
    build_pptx(),
    file_name="digital_transformation_assessment.pptx",
    mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
)
